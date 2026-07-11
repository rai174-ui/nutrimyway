import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # 0 = Title
    # 1 = Title and Content
    # 5 = Title Only
    # 6 = Blank

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "NutriMyWay Portal Overview"
    subtitle.text = "Members and Admin Settings\nDetailed Walkthrough"

    # Slide 2: Members Page Details
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Members Page: Features & Functionalities"
    tf = body_shape.text_frame
    
    p = tf.paragraphs[0]
    p.text = "The Members page is the central hub for managing all members."
    
    p = tf.add_paragraph()
    p.text = "Member Lookup & Onboarding"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Search by Mobile, Email, or Membership No."
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Quick registration for new members directly from the interface."
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Health Records Tracking"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Logs Weight, BMI, Body Fat, Visceral Fat, and Muscle Mass."
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Includes Metabolic Age, BMR, and Resting Heart Rate."
    p.level = 2

    # Slide 3: Members Page Screenshot
    blank_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "Members Page Screenshot"
    
    img_path = r"c:\Users\ABC\projects\nutrimyway\screenshots\member-log.jpg"
    # Centering image
    try:
        # Add image, leaving space for title
        pic = slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(9))
    except Exception as e:
        print(f"Failed to load image {img_path}: {e}")

    # Slide 4: Admin (Settings) Page Details
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Admin Settings: Features & Functionalities"
    tf = body_shape.text_frame
    
    p = tf.paragraphs[0]
    p.text = "Configures global settings and operational parameters for the center."
    
    p = tf.add_paragraph()
    p.text = "Broadcast Settings"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Schedule automated announcements and messages to members."
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Manage retention periods and broadcast history."
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Flavour & Inventory Management"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Configure active operational days (Mon-Sun)."
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Define and manage units of measurement (g, kg, ml, pcs, etc.)."
    p.level = 2

    # Slide 5: Admin (Settings) Page Screenshot
    slide = prs.slides.add_slide(blank_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "Admin Settings Screenshot"
    
    img_path = r"c:\Users\ABC\projects\nutrimyway\screenshots\admin-settings.jpg"
    try:
        pic = slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(9))
    except Exception as e:
        print(f"Failed to load image {img_path}: {e}")

    # Slide 6: Summary
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Summary & Conclusion"
    tf = body_shape.text_frame
    
    p = tf.paragraphs[0]
    p.text = "Streamlined Member Management"
    
    p = tf.add_paragraph()
    p.text = "Provides easy access to health metrics and membership statuses."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Efficient Center Operations"
    
    p = tf.add_paragraph()
    p.text = "Centralizes communication and inventory configurations."
    p.level = 1

    out_file = r"c:\Users\ABC\projects\nutrimyway\NutriMyWay_Walkthrough.pptx"
    prs.save(out_file)
    print(f"Presentation saved to {out_file}")

if __name__ == '__main__':
    create_presentation()
